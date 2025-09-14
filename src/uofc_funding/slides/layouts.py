"""
Slide layout definitions.

This module defines the layout positions and styles for PowerPoint slides.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import TYPE_CHECKING

from pptx.util import Inches

if TYPE_CHECKING:
    from ..config import Config


@dataclass
class Position:
    """
    Represents a position and size on a slide in inches.

    Attributes:
        left: Distance from left edge
        top: Distance from top edge
        width: Width of the element
        height: Height of the element
    """

    left: float
    top: float
    width: float
    height: float

    def to_inches(self) -> tuple[Inches, Inches, Inches, Inches]:
        """Convert all values to pptx Inches objects."""
        return (Inches(self.left), Inches(self.top), Inches(self.width), Inches(self.height))

    def offset(
        self, left: float = 0, top: float = 0, width: float = 0, height: float = 0
    ) -> Position:
        """
        Create a new Position with offsets applied.

        Args:
            left: Offset to add to left position
            top: Offset to add to top position
            width: Offset to add to width
            height: Offset to add to height

        Returns:
            New Position with offsets applied
        """
        return Position(self.left + left, self.top + top, self.width + width, self.height + height)


class SlideLayouts:
    """
    Provides standard layout positions for slides.

    All positions are in inches and configured for 10x7.5" slides (default PowerPoint size).
    """

    def __init__(self, config: Config) -> None:
        """
        Initialize layouts from configuration.

        Args:
            config: Configuration object with slide position settings
        """
        self.config = config
        self._load_positions()

    def _load_positions(self) -> None:
        """Load positions from config or use defaults."""
        positions = self.config.slides.positions

        # Full view chart (centered, large)
        self.full_chart = Position(*positions.get("full_chart", [1.0, 1.0, 8.0, 6.0]))

        # Split view (chart left, table right)
        self.split_chart = Position(*positions.get("split_chart", [0.5, 1.0, 6.2, 6.0]))
        self.split_table = Position(*positions.get("split_table", [6.7, 1.0, 2.8, 6.0]))

        # Dual chart layout (pie left, bar right)
        self.pie_left = Position(*positions.get("pie_left", [0.5, 1.0, 4.5, 5.0]))
        self.bar_right = Position(*positions.get("bar_right", [5.0, 1.0, 4.5, 5.0]))

        # Table positions
        self.table_left = Position(*positions.get("table_left", [0.5, 5.2, 4.5, 2.0]))
        self.table_right = Position(*positions.get("table_right", [5.0, 5.2, 4.5, 2.0]))
        self.centered_table = Position(*positions.get("centered_table", [2.75, 5.2, 4.5, 2.0]))

        # Header/logo position
        self.header_logo = Position(*positions.get("header_logo", [8.0, 6.5, 2.0, 1.0]))

        # Title position
        self.title = Position(*positions.get("title", [0.5, 0.2, 9.0, 0.5]))

    @property
    def slide_width(self) -> float:
        """Get slide width in inches."""
        return self.config.slides.width_inches

    @property
    def slide_height(self) -> float:
        """Get slide height in inches."""
        return self.config.slides.height_inches


# Predefined slide types for common layouts
class SlideType:
    """Enumeration of slide layout types."""

    TITLE = "title"
    FULL_CHART = "full_chart"
    CHART_WITH_TABLE = "chart_with_table"
    DUAL_CHART = "dual_chart"
    DUAL_CHART_WITH_TABLES = "dual_chart_with_tables"
    TABLE_ONLY = "table_only"
