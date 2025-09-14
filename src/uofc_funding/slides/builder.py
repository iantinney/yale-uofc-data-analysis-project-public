"""
PowerPoint presentation builder.

This module handles the creation of PowerPoint presentations from
generated charts and analysis data.
"""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from rich.console import Console

from .layouts import Position, SlideLayouts

if TYPE_CHECKING:
    from ..analyzer import FundingAnalysis
    from ..config import Config

console = Console()


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    r, g, b = tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


def px_to_inches(px: int, dpi: int = 96) -> float:
    """Convert pixels to inches."""
    return px / dpi


class PresentationBuilder:
    """
    Builds PowerPoint presentations from funding analysis data.

    This class handles:
    - Creating the presentation structure
    - Adding title slides with branding
    - Inserting charts and tables
    - Managing slide layouts
    - Scaling images appropriately
    """

    def __init__(self, config: Config, charts_dir: Path, round_number: int) -> None:
        """
        Initialize the presentation builder.

        Args:
            config: Configuration object
            charts_dir: Directory containing generated chart images
            round_number: The funding round number
        """
        self.config = config
        self.charts_dir = Path(charts_dir)
        self.round_number = round_number
        self.layouts = SlideLayouts(config)

        # Create presentation
        self.prs = Presentation()

        # Branding colors
        self.primary_color = hex_to_rgb(config.branding.primary_color)
        self.secondary_color = hex_to_rgb(config.branding.secondary_color)

    def build(self, analysis: FundingAnalysis) -> Presentation:
        """
        Build the complete presentation.

        Args:
            analysis: FundingAnalysis object with computed statistics

        Returns:
            The completed Presentation object
        """
        console.print("[cyan]Building presentation...[/cyan]")

        # Slide 1: Title
        self._add_title_slide()

        # Slide 2: Overview/Rubric
        if analysis.overview_formatted is not None:
            self._add_overview_slide()

        # Slide 3: Top categories (pie + table)
        if analysis.top_categories is not None:
            self._add_categories_slide()

        # Slide 4: Average by category (bar chart)
        if analysis.avg_by_category is not None:
            self._add_avg_category_slide()

        # Slide 5: Age group analysis (pie + bar + tables)
        if analysis.age_group_counts is not None:
            self._add_age_group_slide()

        # Slide 6: Dwight Hall analysis
        if analysis.dwight_hall_counts is not None:
            self._add_dwight_hall_slide()

        # Slide 7: Other funding sources
        if analysis.other_funding_counts is not None:
            self._add_other_funding_slide()

        # Slide 8: Other funding comparison
        if analysis.avg_by_other_funding is not None:
            self._add_other_funding_comparison_slide()

        # Slide 9: Top organizations
        if analysis.top_organizations is not None:
            self._add_top_organizations_slide()

        console.print(f"[green]✓[/green] Created {len(self.prs.slides)} slides")

        return self.prs

    def save(self, output_path: str | Path) -> Path:
        """
        Save the presentation to a file.

        Args:
            output_path: Path for the output file

        Returns:
            Path to the saved file
        """
        output_path = Path(output_path)
        self.prs.save(str(output_path))
        console.print(f"[green]✓[/green] Saved presentation to {output_path}")
        return output_path

    def _add_title_slide(self) -> None:
        """Add the title slide with branding."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.primary_color

        # Add logo if exists
        logo_path = self._find_logo()
        if logo_path:
            slide.shapes.add_picture(str(logo_path), Inches(0.5), Inches(1.5), width=Inches(9))

        # Add title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Round {self.round_number} Transparency Report"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # Add website link
        link_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1.0))
        tf = link_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.config.branding.website_url
        p.font.color.rgb = RGBColor(0, 120, 215)
        p.font.underline = True
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER

    def _add_overview_slide(self) -> None:
        """Add the overview/rubric slide."""
        slide = self._add_slide_with_header(f"Round {self.round_number} Overview")

        chart_path = self.charts_dir / "rubric.png"
        if chart_path.exists():
            self._add_scaled_image(slide, chart_path, self.layouts.full_chart)

    def _add_categories_slide(self) -> None:
        """Add the top categories slide."""
        slide = self._add_slide_with_header("Number of Organizations Within Top 10 Categories")

        # Pie chart on left
        pie_path = self.charts_dir / "top_10_categories_pie.png"
        if pie_path.exists():
            self._add_scaled_image(slide, pie_path, self.layouts.split_chart)

        # Table on right
        csv_path = self.charts_dir / "top_10_categories_table.csv"
        if csv_path.exists():
            self._add_table_from_csv(slide, csv_path, self.layouts.split_table)

    def _add_avg_category_slide(self) -> None:
        """Add the average by category slide."""
        slide = self._add_slide_with_header("Average Amount Requested/Awarded by Category")

        chart_path = self.charts_dir / "avg_requested_by_category_bar.png"
        if chart_path.exists():
            self._add_scaled_image(slide, chart_path, self.layouts.full_chart)

    def _add_age_group_slide(self) -> None:
        """Add the age group analysis slide."""
        slide = self._add_slide_with_header("Amount Requested & Awarded by Organization Age")

        # Pie chart on left
        pie_path = self.charts_dir / "age_group_pie.png"
        if pie_path.exists():
            self._add_scaled_image(slide, pie_path, self.layouts.pie_left)

        # Bar chart on right
        bar_path = self.charts_dir / "avg_req_award_by_age_bar.png"
        if bar_path.exists():
            self._add_scaled_image(slide, bar_path, self.layouts.bar_right)

        # Tables below
        age_csv = self.charts_dir / "age_group_table.csv"
        if age_csv.exists():
            self._add_table_from_csv(
                slide, age_csv, self.layouts.table_left.offset(left=-0.3, top=-0.5)
            )

        avg_csv = self.charts_dir / "avg_req_award_by_age_table.csv"
        if avg_csv.exists():
            self._add_table_from_csv(
                slide, avg_csv, self.layouts.table_right.offset(left=0.3, top=-0.5)
            )

    def _add_dwight_hall_slide(self) -> None:
        """Add the Dwight Hall analysis slide."""
        slide = self._add_slide_with_header("Amount Requested & Awarded by Dwight Hall Group")

        # Pie chart on left
        pie_path = self.charts_dir / "dwight_hall_distribution_pie_chart.png"
        if pie_path.exists():
            self._add_scaled_image(slide, pie_path, self.layouts.pie_left)

        # Bar chart on right
        bar_path = self.charts_dir / "avg_req_award_dwight_hall_bar.png"
        if bar_path.exists():
            self._add_scaled_image(slide, bar_path, self.layouts.bar_right)

        # Table in center below
        table_path = self.charts_dir / "dwight_hall_distribution_table.png"
        if table_path.exists():
            self._add_scaled_image(slide, table_path, self.layouts.centered_table)

    def _add_other_funding_slide(self) -> None:
        """Add the other funding sources slide."""
        slide = self._add_slide_with_header("Other Funding Sources")

        # Pie chart centered
        pie_path = self.charts_dir / "other_funding_sources_pie.png"
        if pie_path.exists():
            centered_pos = self.layouts.split_chart.offset(left=1.5, top=-0.3)
            self._add_scaled_image(slide, pie_path, centered_pos)

        # Table below
        csv_path = self.charts_dir / "other_funding_sources_counts.csv"
        if csv_path.exists():
            self._add_table_from_csv(slide, csv_path, self.layouts.centered_table)

    def _add_other_funding_comparison_slide(self) -> None:
        """Add the other funding comparison slide."""
        slide = self._add_slide_with_header("Amount Requested & Awarded vs Other Funding")

        # Bar chart
        bar_path = self.charts_dir / "avg_req_award_other_funding_bar.png"
        if bar_path.exists():
            self._add_scaled_image(slide, bar_path, self.layouts.split_chart.offset(top=0.7))

        # Table
        csv_path = self.charts_dir / "avg_req_award_other_funding_table.csv"
        if csv_path.exists():
            self._add_table_from_csv(
                slide, csv_path, self.layouts.split_table.offset(top=0.7, height=-2)
            )

    def _add_top_organizations_slide(self) -> None:
        """Add the top organizations slide."""
        slide = self._add_slide_with_header("Top 5 Organizations Summary")

        table_path = self.charts_dir / "top_5_organizations_table.png"
        if table_path.exists():
            self._add_scaled_image(slide, table_path, self.layouts.full_chart)

    def _add_slide_with_header(self, title: str) -> object:
        """
        Add a new slide with title and header logo.

        Args:
            title: Slide title text

        Returns:
            The created slide object
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Title only layout

        # Add title
        self._add_title(slide, title)

        # Add header logo
        self._add_header_logo(slide)

        return slide

    def _add_title(self, slide: object, title: str) -> None:
        """Add a title to the slide."""
        left, top, width, height = self.layouts.title.to_inches()

        title_shape = slide.shapes.add_textbox(left, top, width, height)
        tf = title_shape.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def _add_header_logo(self, slide: object) -> None:
        """Add the header logo to bottom right of slide."""
        logo_path = self._find_logo()
        if not logo_path:
            return

        scale_factor = 0.2

        with Image.open(logo_path) as img:
            img_width = px_to_inches(img.width) * scale_factor
            img_height = px_to_inches(img.height) * scale_factor

        # Position in bottom right
        left = self.layouts.slide_width - img_width - 0.1
        top = self.layouts.slide_height - img_height - 0.1

        slide.shapes.add_picture(
            str(logo_path), Inches(left), Inches(top), Inches(img_width), Inches(img_height)
        )

    def _add_scaled_image(self, slide: object, img_path: Path, position: Position) -> None:
        """
        Add an image to the slide, scaled to fit the position.

        Args:
            slide: The slide object
            img_path: Path to the image file
            position: Target position for the image
        """
        if not img_path.exists():
            console.print(f"[yellow]Warning:[/yellow] Image not found: {img_path}")
            return

        with Image.open(img_path) as img:
            img_width_in = px_to_inches(img.width)
            img_height_in = px_to_inches(img.height)

        # Calculate scale to fit within position while maintaining aspect ratio
        width_ratio = position.width / img_width_in
        height_ratio = position.height / img_height_in
        scale = min(width_ratio, height_ratio, 1.0)  # Don't scale up

        new_width = img_width_in * scale
        new_height = img_height_in * scale

        slide.shapes.add_picture(
            str(img_path),
            Inches(position.left),
            Inches(position.top),
            Inches(new_width),
            Inches(new_height),
        )

    def _add_table_from_csv(self, slide: object, csv_path: Path, position: Position) -> None:
        """
        Add a table to the slide from a CSV file.

        Args:
            slide: The slide object
            csv_path: Path to the CSV file
            position: Position for the table
        """
        if not csv_path.exists():
            console.print(f"[yellow]Warning:[/yellow] CSV not found: {csv_path}")
            return

        df = pd.read_csv(csv_path)
        rows, cols = df.shape

        left, top, width, height = position.to_inches()

        # Create table
        table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

        # Set column widths
        col_width = position.width / cols
        for i in range(cols):
            table.columns[i].width = Inches(col_width)

        # Add header row
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

        # Add data rows
        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table.cell(row_idx + 1, col_idx)
                value = df.iloc[row_idx, col_idx]
                cell.text = str(value) if pd.notna(value) else ""
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER

    def _find_logo(self) -> Path | None:
        """Find the logo file in various locations."""
        possible_paths = [
            Path(self.config.branding.logo_path),
            self.charts_dir.parent / "assets" / "uofcsubtitle.png",
            self.charts_dir.parent / "uofcsubtitle.png",
            Path("uofcsubtitle.png"),
            Path("assets/uofcsubtitle.png"),
        ]

        for path in possible_paths:
            if path.exists():
                return path

        return None
